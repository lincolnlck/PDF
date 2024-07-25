import streamlit as st
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image
import io
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm

def run_phraise_app():
    st.title("Phraise App")
    st.header("이곳은 Phraise App 섹션입니다.")
    # Phraise App의 코드 로직을 여기에 추가


title_align_map = {
    "Left": PP_ALIGN.LEFT,
    "Center": PP_ALIGN.CENTER,
    "Right": PP_ALIGN.RIGHT
}

# 페이지 설정
st.set_page_config(page_title="PPT Lyrics Generator", layout="wide")

# 사용자 정의 CSS 추가
st.markdown(
    """
    <style>
    .main {
        background-color: #e6e6e6;
        padding: 20px;
    }
    .custom-subheader {
        font-size: 24px !important;
        color: #333333;
        font-weight: bold;
        margin-bottom: 5px;
    }
    .stButton>button {
        background-color: #4CAF50;
        color: white;
        border: none;
        padding: 10px 20px;
        text-align: center;
        text-decoration: none;
        display: inline-block;
        font-size: 16px;
        margin: 4px 2px;
        cursor: pointer;
        border-radius: 16px;
    }
    .stRadio label {
        font-size: 20px;
        font-weight: bold;
    }
    .stTextInput input {
        font-size: 16px;
        padding: 10px;
        border: 1px solid #ccc;
        border-radius: 4px;
    }
    .stFileUploader label {
        font-size: 16px;
        font-weight: bold;
    }
    </style>
    """,
    unsafe_allow_html=True,
)

st.title("PPT Lyrics Generator")
st.markdown("🎵 This program automatically generates multiple PPT slides according to the background and text box configuration set by the user. 🎵")

# PPT 사이즈 선택
st.markdown('<p class="custom-subheader">PPT Size</p>', unsafe_allow_html=True)
ppt_size = st.radio("Select PPT Size", ('16:9', '4:3', 'A4 Landscape', 'A4 Portrait'))

# 파일 제목 입력 텍스트 박스 추가
st.markdown('<p class="custom-subheader">Enter the title for your PPT:</p>', unsafe_allow_html=True)
ppt_title = st.text_input("PPT Title", value="", placeholder="My Presentation")

# 배경 이미지 입력
st.markdown('<p class="custom-subheader">Background Upload</p>', unsafe_allow_html=True)
background_image = st.file_uploader("Upload Background Image", type=["jpg", "png"])

# 배경 색상 선택
st.markdown('<p class="custom-subheader">Background Color</p>', unsafe_allow_html=True)
background_color = st.color_picker("Pick Background Color", value='#000000')

# 텍스트 붙여넣기 시 불필요한 줄 바꿈 제거
def clean_text(text):
    return '\n'.join(line for line in text.splitlines() if line.strip() != "")

# 타이틀 설정 추가 버튼
show_title_settings = st.checkbox('Add Title')

# 타이틀 설정 기능
if show_title_settings:
    st.markdown('<p class="custom-subheader">Title Text Box:</p>', unsafe_allow_html=True)
    title_text = st.text_input("Title Text", value="", placeholder="Enter title here...")
    title_text = clean_text(title_text)
    
    col1, col2 = st.columns(2)
    with col1:
        title_font = st.text_input("Font for Title Text", value='서울한강 장체 B', key='title_font')
        title_color = st.color_picker("Font Color for Title Text", value='#FFFFFF', key='title_color')
        title_alignment = st.radio("Alignment for Title Text", ('Left', 'Center', 'Right'), key='title_alignment')
        title_align_map = {'Left': PP_ALIGN.LEFT, 'Center': PP_ALIGN.CENTER, 'Right': PP_ALIGN.RIGHT}
        title_size = st.slider("Font Size for Title Text", 20, 100, 36, key='title_size')
    with col2:
        title_left = st.number_input("Box Left Position (cm) for Title Text", value=1.55, step=0.1, key='title_left')
        title_top = st.number_input("Box Top Position (cm) for Title Text", value=1.66, step=0.1, key='title_top')
        title_width = st.number_input("Box Width (cm) for Title Text", value=30.87, step=0.1, key='title_width')
        title_height = st.number_input("Box Height (cm) for Title Text", value=1.8, step=0.1, key='title_height')

default_lyrics_color = '#FFFFFF'  # 기본 흰색
default_english_lyrics_color = '#FFFF00'  # 기본 밝은 노란색
default_lyrics_font_name = '서울한강 장체 B'
default_english_font_name = '서울한강 장체 B'

# 기본 텍스트 박스 설정을 두 개의 컬럼으로 나누기
col1, col2 = st.columns(2)

with col1:
    st.markdown('<p class="custom-subheader">Text Box 1:</p>', unsafe_allow_html=True)
    text_box_1 = st.text_area("Lyrics Text Box", height=200, placeholder="Enter lyrics here...")
    text_box_1 = clean_text(text_box_1)
    font_1 = st.text_input("Font for Text Box 1", value=default_lyrics_font_name, key='font_1')
    color_1 = st.color_picker("Font Color for Text Box 1", value=default_lyrics_color, key='color_1')
    size_1 = st.slider("Font Size for Text Box 1", 20, 100, 60, key='size_1')
    left_1 = st.number_input("Box Left Position (cm) for Text Box 1", value=0.0, step=0.1, key='left_1')
    top_1 = st.number_input("Box Top Position (cm) for Text Box 1", value=7.0, step=0.1, key='top_1')
    width_1 = st.number_input("Box Width (cm) for Text Box 1", value=33.86, step=0.1, key='width_1')
    height_1 = st.number_input("Box Height (cm) for Text Box 1", value=5.39, step=0.1, key='height_1')

with col2:
    st.markdown('<p class="custom-subheader">Text Box 2:</p>', unsafe_allow_html=True)
    text_box_2 = st.text_area("English Translation Text Box", height=200, placeholder="Enter English translation here...")
    text_box_2 = clean_text(text_box_2)
    font_2 = st.text_input("Font for Text Box 2", value=default_english_font_name, key='font_2')
    color_2 = st.color_picker("Font Color for Text Box 2", value=default_english_lyrics_color, key='color_2')
    size_2 = st.slider("Font Size for Text Box 2", 20, 100, 48, key='size_2')
    left_2 = st.number_input("Box Left Position (cm) for Text Box 2", value=0.0, step=0.1, key='left_2')
    top_2 = st.number_input("Box Top Position (cm) for Text Box 2", value=13.1, step=0.1, key='top_2')
    width_2 = st.number_input("Box Width (cm) for Text Box 2", value=33.86, step=0.1, key='width_2')
    height_2 = st.number_input("Box Height (cm) for Text Box 2", value=3.68, step=0.1, key='height_2')

def hex_to_rgb(hex):
    return tuple(int(hex[i:i+2], 16) for i in (1, 3, 5))

def compress_image(image, max_size_kb, initial_quality=95, min_quality=70):
    img_format = image.format
    output = io.BytesIO()
    
    quality = initial_quality  # 초기 압축 품질 설정
    image.save(output, format=img_format, quality=quality)
    
    while output.tell() > max_size_kb * 1024 and quality > min_quality:
        quality -= 5
        output = io.BytesIO()
        image.save(output, format=img_format, quality=quality)
    
    output.seek(0)
    return output

def create_ppt(text_boxes, fonts, colors, sizes, positions, title, size, bg_image=None, bg_color=None):
    prs = Presentation()
    if size == '16:9':
        prs.slide_width = Cm(33.87)
        prs.slide_height = Cm(19.05)
    elif size == '4:3':
        prs.slide_width = Cm(25.4)
        prs.slide_height = Cm(19.05)
    elif size == 'A4 Landscape':
        prs.slide_width = Cm(29.7)
        prs.slide_height = Cm(21.0)
    else:  # A4 Portrait
        prs.slide_width = Cm(21.0)
        prs.slide_height = Cm(29.7)
    
    lyrics_lines = text_boxes[0].split('\n')
    english_lines = text_boxes[1].split('\n')
    max_lines = max(len(lyrics_lines), len(english_lines))
    
    for i in range(0, max_lines, 2):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        if bg_image:
            slide.shapes.add_picture(bg_image, 0, 0, width=prs.slide_width, height=prs.slide_height)
        else:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*hex_to_rgb(bg_color))  # 사용자 지정 배경 색상 적용
        
        # 가사 텍스트 상자 추가
        text_box_1 = slide.shapes.add_textbox(Cm(positions[0]['left']), Cm(positions[0]['top']), Cm(positions[0]['width']), Cm(positions[0]['height']))
        text_frame_1 = text_box_1.text_frame
        text_frame_1.clear()  # 기본 빈 문단 제거
        text_frame_1.vertical_anchor = MSO_ANCHOR.MIDDLE
        text_frame_1.word_wrap = True
        while text_frame_1.paragraphs:
            text_frame_1._element.remove(text_frame_1.paragraphs[0]._element)

        lyrics_text = '\n'.join(lyrics_lines[i:i+2])  # 두 줄씩 추가
        p1 = text_frame_1.add_paragraph()
        p1.text = lyrics_text
        p1.font.size = Pt(sizes[0])
        p1.font.name = fonts[0]
        p1.font.color.rgb = RGBColor(*hex_to_rgb(colors[0]))
        p1.alignment = PP_ALIGN.CENTER
        
        # 영어 번역 텍스트 상자 추가
        if english_lines:
            text_box_2 = slide.shapes.add_textbox(Cm(positions[1]['left']), Cm(positions[1]['top']), Cm(positions[1]['width']), Cm(positions[1]['height']))
            text_frame_2 = text_box_2.text_frame
            text_frame_2.clear()  # 기본 빈 문단 제거
            text_frame_2.vertical_anchor = MSO_ANCHOR.MIDDLE
            text_frame_2.word_wrap = True
            while text_frame_2.paragraphs:
                text_frame_2._element.remove(text_frame_2.paragraphs[0]._element)

            english_text = '\n'.join(english_lines[i:i+2])  # 두 줄씩 추가
            p2 = text_frame_2.add_paragraph()
            p2.text = english_text
            p2.font.size = Pt(sizes[1])
            p2.font.name = fonts[1]
            p2.font.color.rgb = RGBColor(*hex_to_rgb(colors[1]))
            p2.alignment = PP_ALIGN.CENTER
        
        # 제목 텍스트 상자 추가
        if show_title_settings and title_text:
            title_box = slide.shapes.add_textbox(Cm(title_left), Cm(title_top), Cm(title_width), Cm(title_height))
            text_frame_title = title_box.text_frame
            text_frame_title.clear()  # 기본 빈 문단 제거
            text_frame_title.vertical_anchor = MSO_ANCHOR.MIDDLE
            text_frame_title.word_wrap = True
            while text_frame_title.paragraphs:
                text_frame_title._element.remove(text_frame_title.paragraphs[0]._element)
            p_title = text_frame_title.add_paragraph()
            p_title.text = title_text
            p_title.font.size = Pt(title_size)
            p_title.font.name = title_font
            p_title.font.color.rgb = RGBColor(*hex_to_rgb(title_color))
            p_title.alignment = title_align_map[title_alignment]

    return prs

# 미리보기 기능
def generate_preview(text_boxes, fonts, colors, sizes, positions, bg_image, ppt_size, title_text, title_font, title_color, title_size, title_alignment, title_left, title_top, title_width, title_height, bg_color):
    # 미리보기 화면 크기 설정
    if ppt_size == '16:9':
        fig_size = (33.87, 19.05)
    elif ppt_size == '4:3':
        fig_size = (25.4, 19.05)
    elif ppt_size == 'A4 Landscape':
        fig_size = (29.7, 21.0)
    else:  # A4 Portrait
        fig_size = (21.0, 29.7)

    fig, ax = plt.subplots(figsize=(fig_size[0] / 2.54, fig_size[1] / 2.54))  # 인치를 센티미터로 변환
    
    # 배경 이미지 설정
    if bg_image:
        img = Image.open(bg_image)
        ax.imshow(img, extent=[0, fig_size[0], 0, fig_size[1]])
    else:
        ax.set_facecolor(bg_color)  # 사용자 지정 배경 색상 적용

    ax.set_xlim(0, fig_size[0])
    ax.set_ylim(0, fig_size[1])
    ax.invert_yaxis()
    ax.set_xticks([])
    ax.set_yticks([])

    lyrics_lines = text_boxes[0].split('\n')
    english_lines = text_boxes[1].split('\n')

    def get_font_properties(font_name, size):
        try:
            return fm.FontProperties(fname=fm.findSystemFonts(fontpaths=None, fontext='ttf')[0], size=size)
        except:
            return fm.FontProperties(size=size)

    if len(lyrics_lines) > 0 and len(english_lines) > 0:
        lyrics_text_1 = lyrics_lines[0]
        lyrics_text_2 = lyrics_lines[1] if len(lyrics_lines) > 1 else ""
        english_text_1 = english_lines[0]
        english_text_2 = english_lines[1] if len(english_lines) > 1 else ""

        ax.text(positions[0]['left'] + positions[0]['width']/2, positions[0]['top'] + positions[0]['height']/4, lyrics_text_1,
                fontsize=sizes[0], color=colors[0], ha='center', va='center', fontproperties=get_font_properties(fonts[0], sizes[0]))
        ax.text(positions[0]['left'] + positions[0]['width']/2, positions[0]['top'] + 3*positions[0]['height']/4, lyrics_text_2,
                fontsize=sizes[0], color=colors[0], ha='center', va='center', fontproperties=get_font_properties(fonts[0], sizes[0]))
        ax.text(positions[1]['left'] + positions[1]['width']/2, positions[1]['top'] + positions[1]['height']/4, english_text_1,
                fontsize=sizes[1], color=colors[1], ha='center', va='center', fontproperties=get_font_properties(fonts[1], sizes[1]))
        ax.text(positions[1]['left'] + positions[1]['width']/2, positions[1]['top'] + 3*positions[1]['height']/4, english_text_2,
                fontsize=sizes[1], color=colors[1], ha='center', va='center', fontproperties=get_font_properties(fonts[1], sizes[1]))

    if title_text:
        if title_alignment == 'Left':
            x_pos = title_left
            ha = 'left'
        elif title_alignment == 'Center':
            x_pos = title_left + title_width / 2
            ha = 'center'
        else:  # Right
            x_pos = title_left + title_width
            ha = 'right'

        ax.text(x_pos, title_top + title_height / 2, title_text,
                fontsize=title_size, color=title_color, ha=ha, va='center', fontproperties=get_font_properties(title_font, title_size))

    st.pyplot(fig)


if st.button("Preview"):
    generate_preview(
        [text_box_1, text_box_2],
        [font_1, font_2],
        [color_1, color_2],
        [size_1, size_2],
        [
            {'left': left_1, 'top': top_1, 'width': width_1, 'height': height_1},
            {'left': left_2, 'top': top_2, 'width': width_2, 'height': height_2}
        ],
        background_image,
        ppt_size,
        title_text if show_title_settings else "", 
        title_font if show_title_settings else "Arial",
        title_color if show_title_settings else "#000000",
        title_size if show_title_settings else 24,
        title_alignment if show_title_settings else "Center",
        title_left if show_title_settings else 1.55,
        title_top if show_title_settings else 1.66,
        title_width if show_title_settings else 30.87,
        title_height if show_title_settings else 1.8,
        background_color
    )

if st.button("Generate PPT"):
    if background_image:
        background_image = Image.open(background_image)
        compressed_bg_image = compress_image(background_image, max_size_kb=1024)  # 최대 1MB로 압축
        prs = create_ppt(
            [text_box_1, text_box_2],
            [font_1, font_2],
            [color_1, color_2],
            [size_1, size_2],
            [
                {'left': left_1, 'top': top_1, 'width': width_1, 'height': height_1},
                {'left': left_2, 'top': top_2, 'width': width_2, 'height': height_2}
            ],
            ppt_title, ppt_size, compressed_bg_image, background_color
        )
    else:
        prs = create_ppt(
            [text_box_1, text_box_2],
            [font_1, font_2],
            [color_1, color_2],
            [size_1, size_2],
            [
                {'left': left_1, 'top': top_1, 'width': width_1, 'height': height_1},
                {'left': left_2, 'top': top_2, 'width': width_2, 'height': height_2}
            ],
            ppt_title, ppt_size, None, background_color
        )
    
    ppt_buffer = io.BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    st.download_button("Download PPT", ppt_buffer, file_name=f"{ppt_title}.pptx")


if __name__ == "__main__":
    run_phraise_app()

    